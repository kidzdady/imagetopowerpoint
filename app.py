import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

def create_vertical_layout_ppt(images, output_filename="custom_layout.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Convert pptx dimensions to pixels
    slide_width_px = int(prs.slide_width.pt)
    slide_height_px = int(prs.slide_height.pt)

    # Define image size in pixels
    image_height = slide_height_px // 3

    # Image 1: Top-right
    img1 = images[0].resize((int(slide_width_px * 0.5), int(image_height * 0.95)))
    img_bytes1 = io.BytesIO()
    img1.save(img_bytes1, format='PNG')
    slide.shapes.add_picture(img_bytes1, left=Inches(2.5), top=Inches(0.0))

    # Image 2: Center
    img2 = images[1].resize((int(slide_width_px * 0.6), int(image_height * 0.95)))
    img_bytes2 = io.BytesIO()
    img2.save(img_bytes2, format='PNG')
    slide.shapes.add_picture(img_bytes2, left=Inches(1.2), top=Inches(2.2))

    # Image 3: Bottom
    img3 = images[2].resize((int(slide_width_px * 0.6), int(image_height * 0.95)))
    img_bytes3 = io.BytesIO()
    img3.save(img_bytes3, format='PNG')
    slide.shapes.add_picture(img_bytes3, left=Inches(0.8), top=Inches(4.4))

    prs.save(output_filename)

# --- Streamlit App ---
st.set_page_config(page_title="Custom Image PowerPoint Layout", layout="centered")
st.title("📊 Create Custom Layout PowerPoint Slide")
st.markdown("Upload **3 images** and generate a PowerPoint with them placed like in your design.")

uploaded_files = st.file_uploader("Upload 3 Images", type=["jpg", "jpeg", "png"], accept_multiple_files=True)

if uploaded_files and len(uploaded_files) == 3:
    images = [Image.open(f).convert("RGB") for f in uploaded_files]

    st.subheader("✅ Preview Uploaded Images")
    st.image(images, width=200, caption=["Image 1", "Image 2", "Image 3"])

    if st.button("🛠️ Generate PowerPoint with Custom Layout"):
        pptx_file = "custom_image_layout.pptx"
        create_vertical_layout_ppt(images, pptx_file)

        with open(pptx_file, "rb") as f:
            st.success("🎉 PowerPoint Created!")
            st.download_button("⬇️ Download PowerPoint", f, file_name=pptx_file,
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
else:
    st.info("Please upload exactly 3 images.")
